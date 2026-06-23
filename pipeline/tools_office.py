# Created: 2026-05-20
# Purpose: VEGA Office tools — create/edit/read xlsx/docx/pptx files, LaTeX compile
# Dependencies: sandbox (openpyxl, xlrd, python-docx, python-pptx), subprocess (latex — host)
# Test Status: under validation

from __future__ import annotations

import json
import subprocess
import tempfile
from pathlib import Path
from typing import Any


# ─── Sandbox execution helper ─────────────────────────────────────────────────

def _office_exec(code: str) -> dict:
    """Office 코드 실행 — 호스트 동봉 인터프리터로 직접 실행 (Docker 제거, INT-1870 Phase C).

    xlsx/docx/pptx 생성·편집이 ~/ 에 정상 기록된다. office 작업은 정해진 라이브러리
    호출이라 무한루프 위험이 없고, path_guard + python_exec 가드가 안전을 담당한다.
    """
    from pipeline.tools_code import python_exec
    return python_exec(code, timeout=60)


def _guard_office_paths(args_json: str) -> str | None:
    """ARGS 의 path/src/dst 값을 접근 정책으로 검증. 위반 시 에러 메시지 반환.

    office 출력이 사용자 denylist 폴더나 시크릿 경로에 쓰이는 것을 막는다.
    호스트 직접 실행은 시스템 격리가 없으므로 이 가드가 마지막 방어선."""
    try:
        from pipeline.path_guard import guard_path
        args = json.loads(args_json)
    except Exception:
        return None  # 파싱 불가면 가드 생략(실행부에서 에러 처리)
    for key in ("path", "src", "dst", "dest", "output", "out_path"):
        val = args.get(key)
        if isinstance(val, str) and val:
            try:
                guard_path(val)
            except PermissionError as e:
                return f"[SAFEGUARD] {e}"
    return None


def _sandbox_call(fn_body: str, args_json: str) -> dict:
    """
    fn_body: Python code to run inside the sandbox (receives args via the ARGS variable).
    The function parses whatever JSON is printed to stdout as the return value.
    Execution errors are wrapped as {"error": ...}.
    """
    guard_err = _guard_office_paths(args_json)
    if guard_err:
        return {"error": guard_err}
    code = f"""
import json, shutil, sys
from pathlib import Path

ARGS = json.loads({repr(args_json)})

def _p(path):
    return Path(path).expanduser().resolve()

def _bak(p):
    if p.exists():
        shutil.copy2(p, p.with_suffix(p.suffix + ".bak"))

{fn_body}

print(json.dumps(_result, ensure_ascii=False, default=str))
"""
    result = _office_exec(code)
    if result.get("error"):
        return {"error": result["error"]}
    if result.get("returncode", 0) != 0:
        return {"error": result.get("stderr", "unknown error")[:500]}
    stdout = result.get("stdout", "").strip()
    # The last JSON line in stdout is the return value
    for line in reversed(stdout.splitlines()):
        line = line.strip()
        if line.startswith("{") or line.startswith("["):
            try:
                return json.loads(line)
            except Exception:
                continue
    return {"error": f"Failed to parse result: {stdout[:300]}"}


# ─── xlsx ──────────────────────────────────────────────────────────────────────

def xlsx_create(path: str, sheets: dict[str, list[list]]) -> dict:
    return _sandbox_call("""
import openpyxl
p = _p(ARGS["path"])
p.parent.mkdir(parents=True, exist_ok=True)
_bak(p)
wb = openpyxl.Workbook()
first = True
for name, rows in ARGS["sheets"].items():
    if first:
        ws = wb.active; ws.title = name; first = False
    else:
        ws = wb.create_sheet(name)
    for row in rows:
        ws.append([v if v is not None else "" for v in row])
wb.save(p)
total = sum(len(r) for r in ARGS["sheets"].values())
_result = {"ok": True, "path": str(p), "sheets": list(ARGS["sheets"].keys()), "rows_written": total}
""", json.dumps({"path": path, "sheets": sheets}))


def xlsx_read(path: str, sheet: str | None = None, max_rows: int = 500) -> dict:
    """Read an xlsx file. Defaults to the active sheet if none specified. Also supports xlrd (.xls)."""
    return _sandbox_call("""
from pathlib import Path
p = _p(ARGS["path"])
if not p.exists():
    _result = {"error": f"파일 없음: {p}"}
else:
    ext = p.suffix.lower()
    if ext == ".xls":
        import xlrd
        wb = xlrd.open_workbook(str(p))
        sh_name = ARGS.get("sheet")
        ws = wb.sheet_by_name(sh_name) if sh_name and sh_name in wb.sheet_names() else wb.sheet_by_index(0)
        rows = []
        for r in range(min(ws.nrows, ARGS.get("max_rows", 500))):
            rows.append([ws.cell_value(r, c) for c in range(ws.ncols)])
        _result = {"path": str(p), "sheet": ws.name, "rows": rows, "total_rows": ws.nrows}
    else:
        import openpyxl
        wb = openpyxl.load_workbook(p, read_only=True, data_only=True)
        sh_name = ARGS.get("sheet")
        ws = wb[sh_name] if sh_name and sh_name in wb.sheetnames else wb.active
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= ARGS.get("max_rows", 500): break
            rows.append([v for v in row])
        _result = {"path": str(p), "sheet": ws.title, "rows": rows, "sheets": wb.sheetnames}
        wb.close()
""", json.dumps({"path": path, "sheet": sheet, "max_rows": max_rows}))


def xlsx_merge(output_path: str, sources: list[dict]) -> dict:
    return _sandbox_call("""
import openpyxl
out = _p(ARGS["output_path"])
out.parent.mkdir(parents=True, exist_ok=True)
_bak(out)
wb_out = openpyxl.Workbook()
wb_out.remove(wb_out.active)
log = []
for src in ARGS["sources"]:
    src_path = _p(src["path"])
    if not src_path.exists():
        log.append(f"파일 없음: {src_path}"); continue
    wb_in = openpyxl.load_workbook(src_path, read_only=True, data_only=True)
    sh = src.get("sheet")
    ws_in = wb_in[sh] if sh and sh in wb_in.sheetnames else wb_in.active
    out_name = src.get("as") or ws_in.title
    base = out_name; idx = 1
    while out_name in wb_out.sheetnames:
        out_name = f"{base}_{idx}"; idx += 1
    ws_out = wb_out.create_sheet(out_name)
    skip = src.get("skip_rows", 0); count = 0
    for i, row in enumerate(ws_in.iter_rows(values_only=True)):
        if i < skip: continue
        ws_out.append([v if v is not None else "" for v in row]); count += 1
    wb_in.close()
    log.append(f"✓ {src_path.name}[{ws_in.title}] → [{out_name}] ({count}행)")
wb_out.save(out)
_result = {"ok": True, "path": str(out), "log": log}
""", json.dumps({"output_path": output_path, "sources": sources}))


def xlsx_style(path: str, sheet: str | None, ranges: list[dict]) -> dict:
    return _sandbox_call("""
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
p = _p(ARGS["path"])
if not p.exists():
    _result = {"error": f"파일 없음: {ARGS['path']}"}
else:
    _bak(p)
    wb = openpyxl.load_workbook(p)
    sh = ARGS.get("sheet")
    ws = wb[sh] if sh and sh in wb.sheetnames else wb.active
    applied = 0
    for spec in ARGS["ranges"]:
        rng = spec.get("range", "")
        bold = spec.get("bold", False); bg = spec.get("bg_color")
        fc = spec.get("font_color"); fs = spec.get("font_size"); align = spec.get("align")
        for row in ws[rng]:
            for cell in row:
                if bold or fc or fs:
                    cell.font = Font(bold=bold, color=fc or None, size=fs or None)
                if bg:
                    cell.fill = PatternFill("solid", fgColor=bg)
                if align:
                    cell.alignment = Alignment(horizontal=align)
                applied += 1
    wb.save(p)
    _result = {"ok": True, "cells_styled": applied}
""", json.dumps({"path": path, "sheet": sheet, "ranges": ranges}))


def xlsx_set_formula(path: str, sheet: str | None, cell: str, formula: str) -> dict:
    return _sandbox_call("""
import openpyxl
p = _p(ARGS["path"])
if not p.exists():
    _result = {"error": f"파일 없음: {ARGS['path']}"}
else:
    _bak(p)
    wb = openpyxl.load_workbook(p)
    sh = ARGS.get("sheet")
    ws = wb[sh] if sh and sh in wb.sheetnames else wb.active
    ws[ARGS["cell"]] = ARGS["formula"]
    wb.save(p)
    _result = {"ok": True, "cell": ARGS["cell"], "formula": ARGS["formula"]}
""", json.dumps({"path": path, "sheet": sheet, "cell": cell, "formula": formula}))


# ─── docx ──────────────────────────────────────────────────────────────────────

def docx_read(path: str) -> dict:
    return _sandbox_call("""
from docx import Document
p = _p(ARGS["path"])
if not p.exists():
    _result = {"error": f"파일 없음: {p}"}
else:
    doc = Document(str(p))
    paragraphs = [{"style": para.style.name, "text": para.text}
                  for para in doc.paragraphs if para.text.strip()]
    tables = [{"table_index": i,
               "rows": [[cell.text for cell in row.cells] for row in tbl.rows]}
              for i, tbl in enumerate(doc.tables)]
    _result = {"path": str(p), "paragraphs": paragraphs, "tables": tables,
               "section_count": len(doc.sections)}
""", json.dumps({"path": path}))


def docx_create(path: str, content: list[dict]) -> dict:
    return _sandbox_call("""
from docx import Document
from docx.shared import Pt
p = _p(ARGS["path"])
p.parent.mkdir(parents=True, exist_ok=True)
_bak(p)
doc = Document()
for block in ARGS["content"]:
    t = block.get("type", "paragraph")
    if t == "heading":
        doc.add_heading(block.get("text", ""), level=block.get("level", 1))
    elif t == "paragraph":
        para = doc.add_paragraph(block.get("text", ""))
        if block.get("bold"):
            for run in para.runs: run.bold = True
    elif t == "table":
        rows = block.get("rows", [])
        if rows:
            tbl = doc.add_table(rows=len(rows), cols=len(rows[0]))
            tbl.style = "Table Grid"
            for r, row in enumerate(rows):
                for c, val in enumerate(row):
                    tbl.cell(r, c).text = str(val) if val is not None else ""
    elif t == "pagebreak":
        doc.add_page_break()
doc.save(str(p))
_result = {"ok": True, "path": str(p), "blocks_written": len(ARGS["content"])}
""", json.dumps({"path": path, "content": content}))


def docx_append(path: str, content: list[dict]) -> dict:
    return _sandbox_call("""
from docx import Document
p = _p(ARGS["path"])
if not p.exists():
    _result = {"error": f"파일 없음: {p}"}
else:
    _bak(p)
    doc = Document(str(p))
    for block in ARGS["content"]:
        t = block.get("type", "paragraph")
        if t == "heading":
            doc.add_heading(block.get("text", ""), level=block.get("level", 1))
        elif t == "paragraph":
            doc.add_paragraph(block.get("text", ""))
        elif t == "table":
            rows = block.get("rows", [])
            if rows:
                tbl = doc.add_table(rows=len(rows), cols=len(rows[0]))
                tbl.style = "Table Grid"
                for r, row in enumerate(rows):
                    for c, val in enumerate(row):
                        tbl.cell(r, c).text = str(val) if val is not None else ""
        elif t == "pagebreak":
            doc.add_page_break()
    doc.save(str(p))
    _result = {"ok": True, "path": str(p), "blocks_appended": len(ARGS["content"])}
""", json.dumps({"path": path, "content": content}))


# ─── pptx ──────────────────────────────────────────────────────────────────────

def pptx_read(path: str) -> dict:
    return _sandbox_call("""
from pptx import Presentation
p = _p(ARGS["path"])
if not p.exists():
    _result = {"error": f"파일 없음: {p}"}
else:
    prs = Presentation(str(p))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = [para.text.strip()
                 for shape in slide.shapes if shape.has_text_frame
                 for para in shape.text_frame.paragraphs
                 if para.text.strip()]
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({"slide": i + 1, "texts": texts, "notes": notes})
    _result = {"path": str(p), "slide_count": len(prs.slides), "slides": slides}
""", json.dumps({"path": path}))


def pptx_create(path: str, slides: list[dict]) -> dict:
    return _sandbox_call("""
from pptx import Presentation
from pptx.util import Inches, Pt
p = _p(ARGS["path"])
p.parent.mkdir(parents=True, exist_ok=True)
_bak(p)
prs = Presentation()
blank = prs.slide_layouts[6]
title_body = prs.slide_layouts[1]
for spec in ARGS["slides"]:
    if spec.get("table"):
        slide = prs.slides.add_slide(blank)
        if spec.get("title"):
            txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
            txb.text_frame.text = spec["title"]
            txb.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            txb.text_frame.paragraphs[0].runs[0].font.bold = True
        rows_data = spec["table"]
        rows_n = len(rows_data); cols_n = max(len(r) for r in rows_data)
        tbl = slide.shapes.add_table(rows_n, cols_n, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5 * rows_n)).table
        for r, row in enumerate(rows_data):
            for c, val in enumerate(row):
                tbl.cell(r, c).text = str(val) if val is not None else ""
    else:
        slide = prs.slides.add_slide(title_body)
        if spec.get("title"): slide.shapes.title.text = spec["title"]
        if spec.get("body") and len(slide.placeholders) > 1:
            slide.placeholders[1].text = spec["body"]
    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
prs.save(str(p))
_result = {"ok": True, "path": str(p), "slides_created": len(ARGS["slides"])}
""", json.dumps({"path": path, "slides": slides}))


def pptx_append_slide(path: str, title: str, body: str = "", notes: str = "") -> dict:
    return _sandbox_call("""
from pptx import Presentation
p = _p(ARGS["path"])
if not p.exists():
    _result = {"error": f"파일 없음: {p}"}
else:
    _bak(p)
    prs = Presentation(str(p))
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = ARGS["title"]
    if ARGS.get("body") and len(slide.placeholders) > 1:
        slide.placeholders[1].text = ARGS["body"]
    if ARGS.get("notes"):
        slide.notes_slide.notes_text_frame.text = ARGS["notes"]
    prs.save(str(p))
    _result = {"ok": True, "path": str(p), "slide_added": len(prs.slides)}
""", json.dumps({"path": path, "title": title, "body": body, "notes": notes}))


# ─── LaTeX (runs on host — xelatex is not in the Docker image) ────────────────

def latex_compile(tex_source: str, output_path: str, engine: str = "xelatex") -> dict:
    """
    Compiles a LaTeX source string to PDF using the host macOS TeX Live installation.
    engine: xelatex (default, supports CJK) or pdflatex
    """
    import shutil as _shutil
    out = Path(output_path).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)

    engines = {
        "xelatex": "/Library/TeX/texbin/xelatex",
        "pdflatex": "/Library/TeX/texbin/pdflatex",
    }
    bin_path = engines.get(engine, engines["xelatex"])
    if not Path(bin_path).exists():
        return {"error": f"LaTeX engine not found: {bin_path} (macOS TeX Live required)"}

    with tempfile.TemporaryDirectory() as tmpdir:
        tex_file = Path(tmpdir) / "doc.tex"
        tex_file.write_text(tex_source, encoding="utf-8")
        run_args = [bin_path, "-interaction=nonstopmode", "-output-directory", tmpdir, str(tex_file)]
        r = subprocess.run(run_args, capture_output=True, text=True, timeout=60, cwd=tmpdir)
        if r.returncode == 0:
            subprocess.run(run_args, capture_output=True, text=True, timeout=60, cwd=tmpdir)
        pdf_tmp = Path(tmpdir) / "doc.pdf"
        if pdf_tmp.exists():
            _shutil.copy2(pdf_tmp, out)
            return {"ok": True, "path": str(out), "engine": engine}
        errors = [l for l in r.stdout.splitlines() if l.startswith("!") or "Error" in l]
        return {"error": "compilation failed", "engine_errors": errors[:20], "returncode": r.returncode}


def latex_template(template: str, variables: dict[str, str]) -> str:
    """Substitutes {{variable}} placeholders in a LaTeX template and returns the source. Can be passed directly to latex_compile."""
    result = template
    for key, val in variables.items():
        result = result.replace("{{" + key + "}}", val)
    return result


# ─── PDF 생성 (reportlab — 한글 CID 폰트 내장, Docker·TeX 불필요) ────────────────

_PDF_FONTS_READY = False


def _ensure_pdf_fonts() -> None:
    """한글 CID 폰트 1회 등록(HYGothic=고딕 본문, HYSMyeongJo=명조). 멱등.
    reportlab 내장 Adobe CJK CMap을 쓰므로 TTF 번들이 필요 없다."""
    global _PDF_FONTS_READY
    if _PDF_FONTS_READY:
        return
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    for fn in ("HYGothic-Medium", "HYSMyeongJo-Medium"):
        try:
            pdfmetrics.registerFont(UnicodeCIDFont(fn))
        except Exception:
            pass
    _PDF_FONTS_READY = True


def _md_inline(text: str) -> str:
    """markdown 인라인(굵게/기울임/코드/링크)을 reportlab Paragraph 마크업으로 변환.
    XML escape 를 먼저 한 뒤 마크업 태그를 삽입한다(주입 방지)."""
    import re
    s = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    s = re.sub(r"`([^`]+)`", r'<font face="Courier">\1</font>', s)
    s = re.sub(r"\*\*([^*]+)\*\*", r"<b>\1</b>", s)
    s = re.sub(r"__([^_]+)__", r"<b>\1</b>", s)
    s = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"<i>\1</i>", s)
    s = re.sub(r"\[([^\]]+)\]\((https?://[^)]+)\)", r'<link href="\2" color="blue">\1</link>', s)
    return s


def pdf_create(path: str, content: str, title: str | None = None) -> dict:
    """markdown/일반텍스트 content 를 PDF 로 생성 (한글 지원).

    동봉 reportlab 으로 in-process 렌더 — Docker·TeX 불필요라 비개발자 배포본 자립.
    지원: 제목·헤딩(#~###), 문단, **굵게**/*기울임*/`코드`/[링크], 글머리·번호 목록,
    코드펜스(```), 표(| a | b |), 구분선(---). 코드/인라인코드는 Courier(ASCII).
    """
    try:
        from pipeline.path_guard import guard_path
        guard_path(path)
    except PermissionError as e:
        return {"error": f"[SAFEGUARD] {e}"}
    except Exception:
        pass  # path_guard 로드 실패 시 렌더는 진행(아래 mkdir/쓰기에서 실패 처리)
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Preformatted,
            ListFlowable, ListItem, Table, TableStyle, HRFlowable,
        )
    except Exception as e:
        return {"error": f"reportlab 미설치 — PDF 생성 불가: {e}"}

    import re
    import shutil as _sh

    _ensure_pdf_fonts()
    BODY = "HYGothic-Medium"
    out = Path(path).expanduser()
    try:
        out.parent.mkdir(parents=True, exist_ok=True)
        if out.exists():
            _sh.copy2(out, out.with_suffix(out.suffix + ".bak"))
    except Exception as e:
        return {"error": f"경로 준비 실패: {e}"}

    body = ParagraphStyle("body", fontName=BODY, fontSize=10.5, leading=16, spaceAfter=6)
    h1 = ParagraphStyle("h1", fontName=BODY, fontSize=20, leading=26, spaceBefore=6, spaceAfter=10)
    h2 = ParagraphStyle("h2", fontName=BODY, fontSize=15, leading=21, spaceBefore=6, spaceAfter=8)
    h3 = ParagraphStyle("h3", fontName=BODY, fontSize=12.5, leading=18, spaceBefore=4, spaceAfter=6)
    code_st = ParagraphStyle("code", fontName="Courier", fontSize=9, leading=12,
                             backColor=colors.HexColor("#f2f2f2"), borderPadding=6, spaceAfter=6)

    flow: list = []
    if title:
        flow.append(Paragraph(_md_inline(title), h1))
        flow.append(Spacer(1, 4))

    lines = content.split("\n")
    i = 0
    bullets: list[str] = []

    def _flush_bullets() -> None:
        nonlocal bullets
        if bullets:
            flow.append(ListFlowable(
                [ListItem(Paragraph(_md_inline(b), body)) for b in bullets],
                bulletType="bullet", leftIndent=14))
            bullets = []

    while i < len(lines):
        st = lines[i].strip()
        if st.startswith("```"):                                   # 코드펜스
            _flush_bullets()
            buf, i = [], i + 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                buf.append(lines[i]); i += 1
            i += 1
            flow.append(Preformatted("\n".join(buf) or " ", code_st))
            continue
        if st.startswith("|") and st.endswith("|"):                # 표
            _flush_bullets()
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not all(set(c) <= set("-: ") for c in cells):   # 구분행 스킵
                    rows.append([Paragraph(_md_inline(c), body) for c in cells])
                i += 1
            if rows:
                tbl = Table(rows, hAlign="LEFT")
                tbl.setStyle(TableStyle([
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f2f2f2")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 6), ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ("TOPPADDING", (0, 0), (-1, -1), 3), ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                ]))
                flow.append(tbl); flow.append(Spacer(1, 6))
            continue
        if st.startswith("### "):
            _flush_bullets(); flow.append(Paragraph(_md_inline(st[4:]), h3)); i += 1; continue
        if st.startswith("## "):
            _flush_bullets(); flow.append(Paragraph(_md_inline(st[3:]), h2)); i += 1; continue
        if st.startswith("# "):
            _flush_bullets(); flow.append(Paragraph(_md_inline(st[2:]), h1)); i += 1; continue
        if st in ("---", "***", "___"):
            _flush_bullets()
            flow.append(HRFlowable(width="100%", color=colors.HexColor("#cccccc"),
                                   spaceBefore=4, spaceAfter=8)); i += 1; continue
        m = re.match(r"[-*+]\s+(.*)", st) or re.match(r"\d+\.\s+(.*)", st)
        if m:
            bullets.append(m.group(1)); i += 1; continue
        if not st:
            _flush_bullets(); flow.append(Spacer(1, 4)); i += 1; continue
        _flush_bullets(); flow.append(Paragraph(_md_inline(st), body)); i += 1

    _flush_bullets()
    if not flow:
        flow.append(Paragraph(" ", body))

    pages = {"n": 0}

    def _count(canvas, d):
        pages["n"] += 1

    try:
        doc = SimpleDocTemplate(
            str(out), pagesize=A4,
            leftMargin=20 * mm, rightMargin=20 * mm, topMargin=20 * mm, bottomMargin=20 * mm,
            title=title or out.stem)
        doc.build(flow, onFirstPage=_count, onLaterPages=_count)
    except Exception as e:
        return {"error": f"PDF 렌더 실패: {e}"}
    return {"ok": True, "path": str(out), "pages": pages["n"], "bytes": out.stat().st_size}


# ─── Tool schemas ─────────────────────────────────────────────────────────────

OFFICE_TOOL_SCHEMAS: list[dict] = [
    {
        "type": "function",
        "name": "xlsx_create",
        "description": "새 xlsx 파일 생성 (샌드박스). 시트와 데이터를 한 번에 지정. 파일이 이미 있으면 .bak 백업 후 덮어씀.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "저장 경로 (~/... 또는 절대경로)"},
                "sheets": {
                    "type": "object",
                    "description": "{'시트명': [[행1값1, 행1값2], [행2값1, ...]]} 형태",
                    "additionalProperties": {"type": "array", "items": {"type": "array"}},
                },
            },
            "required": ["path", "sheets"],
        },
    },
    {
        "type": "function",
        "name": "xlsx_read",
        "description": "xlsx/xls 파일 읽기 (샌드박스). .xls는 xlrd, .xlsx는 openpyxl 사용.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "읽을 파일 경로"},
                "sheet": {"type": "string", "description": "시트 이름 (생략 시 활성 시트)"},
                "max_rows": {"type": "integer", "description": "읽을 최대 행 수 (기본 500)", "default": 500},
            },
            "required": ["path"],
        },
    },
    {
        "type": "function",
        "name": "xlsx_merge",
        "description": "여러 xlsx 파일의 시트를 하나의 파일로 합침 (샌드박스).",
        "parameters": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string", "description": "출력 xlsx 경로"},
                "sources": {
                    "type": "array",
                    "description": "[{path, sheet(선택), as(출력시트명, 선택), skip_rows(선택)}]",
                    "items": {"type": "object"},
                },
            },
            "required": ["output_path", "sources"],
        },
    },
    {
        "type": "function",
        "name": "xlsx_style",
        "description": "xlsx 셀 범위에 스타일 적용 (굵기, 배경색, 폰트색, 크기, 정렬) (샌드박스).",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "sheet": {"type": "string", "description": "시트 이름 (생략 시 활성 시트)"},
                "ranges": {
                    "type": "array",
                    "description": "[{range:'A1:D1', bold:true, bg_color:'4472C4', font_color:'FFFFFF', font_size:12, align:'center'}]",
                    "items": {"type": "object"},
                },
            },
            "required": ["path", "ranges"],
        },
    },
    {
        "type": "function",
        "name": "xlsx_set_formula",
        "description": "xlsx 특정 셀에 수식 입력 (샌드박스). 예: formula='=SUM(B2:B10)'",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "sheet": {"type": "string"},
                "cell": {"type": "string", "description": "셀 주소 (예: B12)"},
                "formula": {"type": "string", "description": "수식 문자열 (= 로 시작)"},
            },
            "required": ["path", "cell", "formula"],
        },
    },
    {
        "type": "function",
        "name": "docx_read",
        "description": "Word(.docx) 문서 읽기 (샌드박스). 단락 텍스트, 표, 스타일 반환.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "docx 파일 경로"},
            },
            "required": ["path"],
        },
    },
    {
        "type": "function",
        "name": "docx_create",
        "description": "새 Word 문서 생성 (샌드박스). content 블록 리스트로 구성.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "content": {
                    "type": "array",
                    "description": "[{type:'heading'|'paragraph'|'table'|'pagebreak', text, level, bold, rows}]",
                    "items": {"type": "object"},
                },
            },
            "required": ["path", "content"],
        },
    },
    {
        "type": "function",
        "name": "docx_append",
        "description": "기존 Word 문서에 내용 추가 (샌드박스). content 형식은 docx_create와 동일.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "content": {"type": "array", "items": {"type": "object"}},
            },
            "required": ["path", "content"],
        },
    },
    {
        "type": "function",
        "name": "pptx_read",
        "description": "PowerPoint(.pptx) 파일 읽기 (샌드박스). 슬라이드별 텍스트·메모 반환.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
            },
            "required": ["path"],
        },
    },
    {
        "type": "function",
        "name": "pptx_create",
        "description": "새 PowerPoint 파일 생성 (샌드박스). 슬라이드 리스트로 구성.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "slides": {
                    "type": "array",
                    "description": "[{title, body, table:[[]], notes}]",
                    "items": {"type": "object"},
                },
            },
            "required": ["path", "slides"],
        },
    },
    {
        "type": "function",
        "name": "pptx_append_slide",
        "description": "기존 PowerPoint에 슬라이드 추가 (샌드박스).",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "title": {"type": "string"},
                "body": {"type": "string", "default": ""},
                "notes": {"type": "string", "default": ""},
            },
            "required": ["path", "title"],
        },
    },
    {
        "type": "function",
        "name": "latex_compile",
        "description": "LaTeX 소스를 컴파일해 PDF 생성 (호스트 TeX Live 사용, 한국어는 xelatex 권장).",
        "parameters": {
            "type": "object",
            "properties": {
                "tex_source": {"type": "string", "description": "컴파일할 LaTeX 소스 전체"},
                "output_path": {"type": "string", "description": "저장할 PDF 경로"},
                "engine": {
                    "type": "string",
                    "enum": ["xelatex", "pdflatex"],
                    "default": "xelatex",
                    "description": "LaTeX 엔진 (기본: xelatex)",
                },
            },
            "required": ["tex_source", "output_path"],
        },
    },
    {
        "type": "function",
        "name": "latex_template",
        "description": "LaTeX 템플릿의 {{변수}} 플레이스홀더를 치환해 소스 문자열 반환. latex_compile에 바로 넘길 수 있음.",
        "parameters": {
            "type": "object",
            "properties": {
                "template": {"type": "string", "description": "{{변수}} 포함 LaTeX 소스"},
                "variables": {
                    "type": "object",
                    "description": "{'변수명': '치환값'} 딕셔너리",
                    "additionalProperties": {"type": "string"},
                },
            },
            "required": ["template", "variables"],
        },
    },
    {
        "type": "function",
        "name": "pdf_create",
        "description": "markdown/텍스트 내용을 PDF로 생성 (한글 지원, Docker·TeX 불필요). 제목·헤딩·**굵게**/*기울임*/`코드`·목록·표·코드블록·구분선 지원. LLM이 만든 보고서·문서를 PDF로 저장할 때 사용. (고품질 조판이 필요하면 latex_compile.)",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "저장할 PDF 경로 (~/... 또는 절대경로, .pdf)"},
                "content": {"type": "string", "description": "PDF 본문 (markdown 또는 일반 텍스트)"},
                "title": {"type": "string", "description": "문서 제목 (선택 — 첫 페이지 상단 + PDF 메타데이터)"},
            },
            "required": ["path", "content"],
        },
    },
]

OFFICE_TOOL_FUNCTIONS: dict[str, Any] = {
    "xlsx_create":       xlsx_create,
    "xlsx_read":         xlsx_read,
    "xlsx_merge":        xlsx_merge,
    "xlsx_style":        xlsx_style,
    "xlsx_set_formula":  xlsx_set_formula,
    "docx_read":         docx_read,
    "docx_create":       docx_create,
    "docx_append":       docx_append,
    "pptx_read":         pptx_read,
    "pptx_create":       pptx_create,
    "pptx_append_slide": pptx_append_slide,
    "latex_compile":     latex_compile,
    "latex_template":    latex_template,
    "pdf_create":        pdf_create,
}
